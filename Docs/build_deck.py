#!/usr/bin/env python3
"""Builds the FBLA Nationals presentation for Clock-In."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette -------------------------------------------------------------
NAVY      = RGBColor(0x0E, 0x1A, 0x2B)   # deep background
NAVY2     = RGBColor(0x16, 0x27, 0x3D)   # panel
CYAN      = RGBColor(0x2D, 0xD4, 0xEE)   # primary accent (packets)
BLUE      = RGBColor(0x38, 0xBD, 0xF8)
AMBER     = RGBColor(0xF5, 0xC2, 0x42)   # job/clock accent
GREEN     = RGBColor(0x34, 0xD3, 0x99)
LIGHT     = RGBColor(0xF1, 0xF5, 0xF9)   # primary text on dark
MUTED     = RGBColor(0x9F, 0xB3, 0xC8)   # secondary text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=NAVY):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, b, it) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = b; r.font.italic = it
            r.font.name = "Calibri"
    return tb

def R(t, sz, col=LIGHT, b=False, it=False):  # run helper
    return (t, sz, col, b, it)

def footer(s, n, section=""):
    txt(s, Inches(0.5), Inches(7.02), Inches(6), Inches(0.35),
        [[R("CLOCK-IN", 9, MUTED, True), R("   |   FBLA National Leadership Conference 2026", 9, MUTED)]],
        anchor=MSO_ANCHOR.MIDDLE)
    if section:
        txt(s, Inches(7), Inches(7.02), Inches(5.3), Inches(0.35),
            [[R(section.upper(), 9, MUTED)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(12.5), Inches(7.02), Inches(0.6), Inches(0.35),
        [[R(str(n), 9, CYAN, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def header(s, kicker, title, accent=CYAN):
    rect(s, Inches(0.5), Inches(0.55), Inches(0.14), Inches(0.95), accent)
    txt(s, Inches(0.78), Inches(0.5), Inches(11), Inches(0.4),
        [[R(kicker.upper(), 13, accent, True)]])
    txt(s, Inches(0.78), Inches(0.85), Inches(11.8), Inches(0.85),
        [[R(title, 30, WHITE, True)]])

def bullet(text, sz=15, col=LIGHT, bold_lead=None):
    runs = []
    if bold_lead:
        runs.append(R(bold_lead, sz, CYAN, True))
        runs.append(R(text, sz, col))
    else:
        runs.append(R(text, sz, col))
    return runs

def card(s, x, y, w, h, accent, title, lines, title_sz=15, body_sz=12):
    rect(s, x, y, w, h, NAVY2)
    rect(s, x, y, w, Inches(0.09), accent)
    body = [[R(title, title_sz, WHITE, True)]]
    for ln in lines:
        body.append([R("•  ", body_sz, accent, True), R(ln, body_sz, MUTED)])
    txt(s, x+Inches(0.22), y+Inches(0.22), w-Inches(0.4), h-Inches(0.4),
        body, space_after=5, line_spacing=1.05)

# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = slide(); bg(s)
rect(s, 0, 0, SW, Inches(0.18), CYAN)
rect(s, 0, Inches(7.32), SW, Inches(0.18), CYAN)
# big clock-in motif
rect(s, Inches(9.7), Inches(1.3), Inches(3.0), Inches(3.0), NAVY2, shape=MSO_SHAPE.OVAL)
rect(s, Inches(10.25), Inches(1.85), Inches(1.9), Inches(1.9), NAVY, shape=MSO_SHAPE.OVAL)
txt(s, Inches(9.7), Inches(2.45), Inches(3.0), Inches(0.7),
    [[R("9 to 5", 24, CYAN, True)]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(1.7), Inches(8.5), Inches(0.5),
    [[R("PARALLEL PLAY  PRESENTS", 14, CYAN, True)]])
txt(s, Inches(0.9), Inches(2.15), Inches(9), Inches(1.6),
    [[R("CLOCK-IN", 76, WHITE, True)]])
txt(s, Inches(0.95), Inches(3.55), Inches(8.5), Inches(0.6),
    [[R("Experience real careers through multiplayer mini-games.", 19, LIGHT, it=True)]])
rect(s, Inches(0.95), Inches(4.35), Inches(4.6), Inches(0.04), CYAN)
txt(s, Inches(0.95), Inches(4.55), Inches(11), Inches(0.45),
    [[R("FBLA National Leadership Conference 2026", 16, BLUE, True)]])
txt(s, Inches(0.95), Inches(4.95), Inches(11), Inches(0.45),
    [[R("Computer Game & Simulation Programming", 14, MUTED)]])
txt(s, Inches(0.95), Inches(5.75), Inches(11), Inches(0.9),
    [[R("CREATED BY", 12, CYAN, True)],
     [R("Jackson Krizan   ·   Blayne Campbell   ·   Jakob Escueta", 17, LIGHT, True)],
     [R("Built in Unity with C#  ·  Photon PUN 2 multiplayer  ·  Firebase cloud", 12, MUTED)]],
    space_after=4)

# =========================================================================
# SLIDE 2 — TABLE OF CONTENTS
# =========================================================================
s = slide(); bg(s)
header(s, "Roadmap", "Table of Contents")
items = [
    ("01", "Introduction", "The idea behind Clock-In", CYAN),
    ("02", "The Game", "Themes, goals & core concepts", CYAN),
    ("03", "The Jobs", "Four careers, one world", AMBER),
    ("04", "Gameplay Spotlight", "Inside the networking mini-game", AMBER),
    ("05", "Multiplayer", "Real-time play with Photon", BLUE),
    ("06", "Live Scoreboard", "Competing in real time", BLUE),
    ("07", "Accounts & Leaderboard", "Cloud progression with Firebase", GREEN),
    ("08", "Tech & Process", "Our stack, tools & teamwork", GREEN),
    ("09", "New Since State", "How far we've come", CYAN),
    ("10", "What's Next", "The road ahead", CYAN),
]
x0, y0 = Inches(0.78), Inches(1.95)
cw, ch = Inches(5.85), Inches(0.92)
for i, (num, t, sub, acc) in enumerate(items):
    col = i // 5; row = i % 5
    x = x0 + col*(cw + Inches(0.35))
    y = y0 + row*(ch + Inches(0.13))
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, Inches(0.09), ch, acc)
    txt(s, x+Inches(0.28), y, Inches(1.0), ch,
        [[R(num, 26, acc, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(1.25), y+Inches(0.13), cw-Inches(1.4), Inches(0.7),
        [[R(t, 16, WHITE, True)], [R(sub, 11, MUTED)]], space_after=2)
footer(s, 2, "Roadmap")

# =========================================================================
# SLIDE 3 — INTRODUCTION
# =========================================================================
s = slide(); bg(s)
header(s, "01  ·  Introduction", "What is Clock-In?")
txt(s, Inches(0.78), Inches(1.95), Inches(7.1), Inches(3.0),
    [[R("Clock-In is a 3D multiplayer game that lets players ", 17, LIGHT),
      R("try on real-world careers", 17, CYAN, True),
      R(" — without leaving the classroom.", 17, LIGHT)],
     [R("", 8, LIGHT)],
     [R("Each job is a hands-on mini-game built around what that career actually involves: routing data on a network, repairing a car, working a kitchen, or responding to a fire. Players jump into a shared world, take on a role, and learn by doing.", 15, MUTED)],
     [R("", 8, LIGHT)],
     [R("Our goal: make career exploration ", 15, MUTED),
      R("fun, social, and interactive", 15, GREEN, True),
      R(" for students deciding what they want to be.", 15, MUTED)]],
    space_after=8, line_spacing=1.12)
# stat panel
panel_x = Inches(8.3)
rect(s, panel_x, Inches(1.95), Inches(4.25), Inches(4.55), NAVY2)
rect(s, panel_x, Inches(1.95), Inches(4.25), Inches(0.09), CYAN)
stats = [("4", "career mini-games", CYAN),
         ("∞", "players, online together", BLUE),
         ("1", "shared multiplayer world", AMBER),
         ("100%", "learn-by-doing", GREEN)]
sy = Inches(2.25)
for val, lab, acc in stats:
    txt(s, panel_x+Inches(0.35), sy, Inches(3.6), Inches(1.0),
        [[R(val, 40, acc, True)], [R(lab.upper(), 12, LIGHT, True)]], space_after=0)
    sy += Inches(1.02)
footer(s, 3, "Introduction")

# =========================================================================
# SLIDE 4 — THE GAME (themes / goals / concepts)
# =========================================================================
s = slide(); bg(s)
header(s, "02  ·  The Game", "Themes, Goals & Concepts")
card(s, Inches(0.78), Inches(1.95), Inches(3.75), Inches(2.15), CYAN, "THEMES",
     ["Real careers, simulated", "Learning through play", "Cooperative & competitive",
      "Accessible to any student"], body_sz=13)
card(s, Inches(4.79), Inches(1.95), Inches(3.75), Inches(2.15), AMBER, "GOALS",
     ["Visualize what careers look like", "Create a fun interactive experience",
      "Spark interest in skilled trades & tech", "Make learning feel like a game"], body_sz=13)
card(s, Inches(8.80), Inches(1.95), Inches(3.75), Inches(2.15), GREEN, "KEY CONCEPTS",
     ["Short, focused job tasks", "Career-quest progression", "Live multiplayer scoring",
      "Persistent cloud accounts"], body_sz=13)
# strip
rect(s, Inches(0.78), Inches(4.45), Inches(11.77), Inches(1.65), NAVY2)
rect(s, Inches(0.78), Inches(4.45), Inches(0.09), Inches(1.65), BLUE)
txt(s, Inches(1.1), Inches(4.65), Inches(11.2), Inches(1.3),
    [[R("The pitch in one line", 14, BLUE, True)],
     [R("“What do you want to be when you grow up?”  Clock-In turns that question into a game — drop into a job, learn the ropes with friends, climb the leaderboard, and try the next career.", 16, LIGHT, it=True)]],
    space_after=6, line_spacing=1.1)
footer(s, 4, "The Game")

# =========================================================================
# SLIDE 5 — THE JOBS
# =========================================================================
s = slide(); bg(s)
header(s, "03  ·  The Jobs", "Four Careers, One World")
jobs = [
    ("NETWORK ENGINEER", CYAN, "Route data packets across a shifting grid of tiles to their destination.",
     "Teaches: packets, routing & timing", "FEATURED"),
    ("AUTO MECHANIC", AMBER, "Grab, rotate and fit the right parts to repair and rebuild a car.",
     "Teaches: assembly & precision work", ""),
    ("CHEF", GREEN, "Work the kitchen line, handling ingredients under pressure.",
     "Teaches: food-service workflow", "IN PROGRESS"),
    ("FIREFIGHTER", RGBColor(0xEF,0x6F,0x6F), "Operate gear and respond to emergencies to save lives.",
     "Teaches: emergency response", "IN PROGRESS"),
]
x0 = Inches(0.78); cw = Inches(2.86); gap = Inches(0.18); y = Inches(1.95); ch = Inches(4.15)
for i, (name, acc, desc, teach, tag) in enumerate(jobs):
    x = x0 + i*(cw+gap)
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, cw, Inches(0.85), acc)
    txt(s, x+Inches(0.2), y+Inches(0.16), cw-Inches(0.4), Inches(0.6),
        [[R(name, 14, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.22), y+Inches(1.05), cw-Inches(0.44), Inches(1.8),
        [[R(desc, 13, LIGHT)]], line_spacing=1.1)
    txt(s, x+Inches(0.22), y+Inches(2.95), cw-Inches(0.44), Inches(0.8),
        [[R(teach, 11, acc, True)]], line_spacing=1.05)
    if tag:
        tcol = GREEN if tag == "FEATURED" else MUTED
        txt(s, x+Inches(0.22), y+Inches(3.7), cw-Inches(0.44), Inches(0.35),
            [[R("● ", 10, tcol), R(tag, 10, tcol, True)]])
footer(s, 5, "The Jobs")

# =========================================================================
# SLIDE 6 — GAMEPLAY SPOTLIGHT: NETWORKING
# =========================================================================
s = slide(); bg(s)
header(s, "04  ·  Gameplay Spotlight", "Inside the Networking Mini-Game")
txt(s, Inches(0.78), Inches(1.85), Inches(7.0), Inches(0.5),
    [[R("Our deepest job — it turns abstract networking ideas into something you can run across.", 14, MUTED, it=True)]],
    line_spacing=1.05)
steps = [
    ("1", "PICK UP THE PACKET", "Cyan-glowing packets appear in the level — grab one to carry it."),
    ("2", "CROSS THE LIVE GRID", "The floor is made of tiles that fade in and out on 3–10s cycles. Time your route."),
    ("3", "DELIVER TO THE GOAL", "Reach the goal zone — the packet turns green and you score points."),
    ("4", "SYNCED FOR EVERYONE", "Packet state syncs over Photon, so every player sees the same delivery live."),
]
y = Inches(2.45)
for num, t, d in steps:
    rect(s, Inches(0.78), y, Inches(7.05), Inches(0.86), NAVY2)
    rect(s, Inches(0.78), y, Inches(0.7), Inches(0.86), CYAN)
    txt(s, Inches(0.78), y, Inches(0.7), Inches(0.86), [[R(num, 26, NAVY, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.65), y+Inches(0.12), Inches(6.0), Inches(0.7),
        [[R(t, 13, WHITE, True)], [R(d, 11, MUTED)]], space_after=1, line_spacing=1.0)
    y += Inches(0.97)
# concept panel
px = Inches(8.15)
rect(s, px, Inches(2.45), Inches(4.4), Inches(3.78), NAVY2)
rect(s, px, Inches(2.45), Inches(4.4), Inches(0.09), GREEN)
txt(s, px+Inches(0.3), Inches(2.65), Inches(3.9), Inches(3.5),
    [[R("WHAT IT TEACHES", 13, GREEN, True)],
     [R("• Packets", 14, CYAN, True), R(" — data travels in pieces", 13, MUTED)],
     [R("• Routing", 14, CYAN, True), R(" — finding a path to the destination", 13, MUTED)],
     [R("• Timing & latency", 14, CYAN, True), R(" — networks change moment to moment", 13, MUTED)],
     [R("• Synchronization", 14, CYAN, True), R(" — keeping every node in agreement", 13, MUTED)],
     [R("", 8, LIGHT)],
     [R("Built from real scripts:", 11, LIGHT, True)],
     [R("Packet · PacketGoal · Tiles · ScoreManager", 11, MUTED, it=True)]],
    space_after=8, line_spacing=1.1)
footer(s, 6, "Gameplay Spotlight")

# =========================================================================
# SLIDE 7 — MULTIPLAYER
# =========================================================================
s = slide(); bg(s)
header(s, "05  ·  Multiplayer", "Playing Together with Photon PUN 2")
txt(s, Inches(0.78), Inches(1.85), Inches(11.6), Inches(0.6),
    [[R("Clock-In is multiplayer to its core. Players connect to a shared lobby, join a room, and drop into the same job together — in real time.", 15, LIGHT)]],
    line_spacing=1.1)
flow = [
    ("CONNECT", "Players join the Photon master server and lobby", CYAN),
    ("ROOM", "Host creates or joins a room; everyone sees the player list", BLUE),
    ("LAUNCH", "Host picks a job and loads the level for the whole room", AMBER),
    ("SPAWN", "Each player spawns the right character for that scene", GREEN),
    ("PLAY", "Movement & scores sync live across all players", CYAN),
]
x0 = Inches(0.78); cw = Inches(2.18); gap = Inches(0.16); y = Inches(2.95); ch = Inches(1.9)
for i, (t, d, acc) in enumerate(flow):
    x = x0 + i*(cw+gap)
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, cw, Inches(0.09), acc)
    txt(s, x+Inches(0.18), y+Inches(0.22), cw-Inches(0.36), Inches(0.5),
        [[R(t, 15, acc, True)]])
    txt(s, x+Inches(0.18), y+Inches(0.72), cw-Inches(0.36), Inches(1.1),
        [[R(d, 12, MUTED)]], line_spacing=1.05)
    if i < len(flow)-1:
        txt(s, x+cw-Inches(0.02), y, gap+Inches(0.1), ch,
            [[R("›", 22, CYAN, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(0.78), Inches(5.25), Inches(11.77), Inches(1.0), NAVY2)
rect(s, Inches(0.78), Inches(5.25), Inches(0.09), Inches(1.0), BLUE)
txt(s, Inches(1.1), Inches(5.4), Inches(11.2), Inches(0.75),
    [[R("Under the hood:  ", 13, BLUE, True),
      R("a persistent RoomManager keeps players in sync across scene loads, a networked PlayerManager spawns the correct character per job, and host-controlled level loading keeps the whole room on the same map.", 13, MUTED)]],
    line_spacing=1.08)
footer(s, 7, "Multiplayer")

# =========================================================================
# SLIDE 8 — LIVE SCOREBOARD
# =========================================================================
s = slide(); bg(s)
header(s, "06  ·  Live Scoreboard", "Competing in Real Time")
txt(s, Inches(0.78), Inches(1.9), Inches(6.6), Inches(1.0),
    [[R("Every delivery, repair, and task earns points. A live scoreboard shows where everyone stands — instantly, for all players in the room.", 15, LIGHT)]],
    line_spacing=1.12)
feats = [
    ("Backend-agnostic core", "A single ScoreManager owns the score and fires events — gameplay never talks to the network directly."),
    ("Synced over Photon", "ScoreNetworkSync pushes each score into Photon properties, so all players see updates live."),
    ("Toggle anytime", "Press a key to pop the full scoreboard, sorted highest-first, mid-match."),
    ("Clean HUD", "A lightweight on-screen score display keeps your own total always visible."),
]
y = Inches(3.0)
for t, d in feats:
    rect(s, Inches(0.78), y, Inches(6.55), Inches(0.74), NAVY2)
    rect(s, Inches(0.78), y, Inches(0.09), Inches(0.74), CYAN)
    txt(s, Inches(1.05), y+Inches(0.09), Inches(6.1), Inches(0.6),
        [[R(t, 13, WHITE, True)], [R(d, 11, MUTED)]], space_after=1, line_spacing=1.0)
    y += Inches(0.82)
# mock scoreboard
mx = Inches(7.75)
rect(s, mx, Inches(1.9), Inches(4.8), Inches(4.55), NAVY2)
rect(s, mx, Inches(1.9), Inches(4.8), Inches(0.6), CYAN)
txt(s, mx+Inches(0.3), Inches(1.9), Inches(4.2), Inches(0.6),
    [[R("LIVE SCOREBOARD", 15, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
rows = [("1", "Jackson", "120", AMBER), ("2", "Blayne", "90", LIGHT),
        ("3", "Jakob", "80", LIGHT), ("4", "Guest_417", "40", MUTED)]
ry = Inches(2.75)
for rank, name, sc, col in rows:
    rect(s, mx+Inches(0.3), ry, Inches(4.2), Inches(0.78), NAVY)
    txt(s, mx+Inches(0.5), ry, Inches(0.6), Inches(0.78), [[R(rank, 22, CYAN, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, mx+Inches(1.15), ry, Inches(2.3), Inches(0.78), [[R(name, 16, col, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, mx+Inches(3.2), ry, Inches(1.1), Inches(0.78), [[R(sc, 20, CYAN, True)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ry += Inches(0.88)
footer(s, 8, "Live Scoreboard")

# =========================================================================
# SLIDE 9 — ACCOUNTS & LEADERBOARD
# =========================================================================
s = slide(); bg(s)
header(s, "07  ·  Accounts & Leaderboard", "Cloud Progression with Firebase", GREEN)
txt(s, Inches(0.78), Inches(1.9), Inches(11.6), Inches(0.6),
    [[R("New for Nationals: ", 15, GREEN, True),
      R("players can create an account, sign in, and have their best scores saved to a global leaderboard in the cloud.", 15, LIGHT)]],
    line_spacing=1.1)
card(s, Inches(0.78), Inches(2.75), Inches(3.75), Inches(3.5), GREEN, "ACCOUNTS",
     ["Sign up with email & display name", "Secure sign-in via Firebase Auth",
      "Continue as Guest — no account needed", "Session persists across scenes",
      "One-tap sign-out"], body_sz=12)
card(s, Inches(4.79), Inches(2.75), Inches(3.75), Inches(3.5), CYAN, "GLOBAL LEADERBOARD",
     ["Personal-best scores saved to cloud", "Firebase Realtime Database",
      "Ranked top-N global standings", "Guests stay private (never written)",
      "Tracks progress across sessions"], body_sz=12)
card(s, Inches(8.80), Inches(2.75), Inches(3.75), Inches(3.5), BLUE, "WHY IT MATTERS",
     ["Turns one match into a season", "Rewards coming back & improving",
      "Casual OR committed — your choice", "Shows real backend engineering",
      "Scales to any number of players"], body_sz=12)
footer(s, 9, "Accounts & Leaderboard")

# =========================================================================
# SLIDE 10 — TECH STACK
# =========================================================================
s = slide(); bg(s)
header(s, "08  ·  Technology", "Our Tech Stack")
tech = [
    ("UNITY 2022.3.22f1", "3D game engine & editor", CYAN),
    ("C#", "All gameplay & systems code", BLUE),
    ("PHOTON PUN 2", "Real-time multiplayer networking", AMBER),
    ("FIREBASE", "Cloud auth & leaderboard database", GREEN),
    ("TEXTMESH PRO", "Crisp in-game UI & text", CYAN),
    ("GIT + GITHUB", "Version control & team collaboration", BLUE),
]
x0 = Inches(0.78); cw = Inches(3.82); gap = Inches(0.16); ch = Inches(1.55)
y0 = Inches(2.1)
for i, (t, d, acc) in enumerate(tech):
    col = i % 3; row = i // 3
    x = x0 + col*(cw+gap); y = y0 + row*(ch+Inches(0.2))
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, Inches(0.09), ch, acc)
    txt(s, x+Inches(0.3), y+Inches(0.28), cw-Inches(0.5), Inches(1.1),
        [[R(t, 18, WHITE, True)], [R(d, 12, MUTED)]], space_after=6)
rect(s, Inches(0.78), Inches(5.65), Inches(11.77), Inches(0.62), NAVY2)
txt(s, Inches(1.05), Inches(5.65), Inches(11.3), Inches(0.62),
    [[R("Architecture note:  ", 12, CYAN, True),
      R("systems are decoupled — the scoring core has no idea whether it's talking to Photon or Firebase, so we can extend or swap pieces without breaking gameplay.", 12, MUTED)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, 10, "Technology")

# =========================================================================
# SLIDE 11 — PROCESS / TEAM
# =========================================================================
s = slide(); bg(s)
header(s, "08  ·  Process", "How We Built It")
card(s, Inches(0.78), Inches(2.0), Inches(3.75), Inches(2.05), CYAN, "JACKSON KRIZAN",
     ["Networking mini-game", "Scoring & multiplayer sync", "Firebase accounts & leaderboard"], body_sz=12)
card(s, Inches(4.79), Inches(2.0), Inches(3.75), Inches(2.05), AMBER, "BLAYNE CAMPBELL",
     ["Game design & concept", "Job themes & content", "Playtesting & feedback"], body_sz=12)
card(s, Inches(8.80), Inches(2.0), Inches(3.75), Inches(2.05), GREEN, "JAKOB ESCUETA",
     ["Programming support", "User interface & menus", "Assets & integration"], body_sz=12)
rect(s, Inches(0.78), Inches(4.4), Inches(11.77), Inches(1.85), NAVY2)
rect(s, Inches(0.78), Inches(4.4), Inches(0.09), Inches(1.85), BLUE)
txt(s, Inches(1.1), Inches(4.6), Inches(11.2), Inches(1.55),
    [[R("Our workflow", 14, BLUE, True)],
     [R("We work in GitHub branches and pull requests — each feature is built in isolation, reviewed, then merged. Tickets track every task. Unity lets us iterate on gameplay fast, while clean code organization (separate folders per job, shared core systems) keeps a four-job game manageable for a three-person team.", 14, LIGHT)]],
    space_after=8, line_spacing=1.12)
footer(s, 11, "Process")

# =========================================================================
# SLIDE 12 — NEW SINCE STATE
# =========================================================================
s = slide(); bg(s)
header(s, "09  ·  Growth", "What's New Since State")
txt(s, Inches(0.78), Inches(1.85), Inches(11.6), Inches(0.5),
    [[R("We didn't just polish — we added entire systems that take Clock-In from a demo to a real, replayable game.", 15, LIGHT, it=True)]],
    line_spacing=1.05)
news = [
    ("CLOUD ACCOUNTS", "Brand-new Firebase sign-up / sign-in with guest mode", GREEN),
    ("GLOBAL LEADERBOARD", "Persistent personal bests ranked across all players", CYAN),
    ("LIVE MULTIPLAYER SCORING", "Real-time scoreboard synced over Photon", BLUE),
    ("NETWORKING MINI-GAME", "Full packet-delivery game with dynamic tiles", AMBER),
    ("MATURE MULTIPLAYER", "Robust rooms, spawning & scene handling", BLUE),
    ("MORE POLISH", "Visual feedback, HUD, cleaner architecture", CYAN),
]
x0 = Inches(0.78); cw = Inches(3.82); gap = Inches(0.16); ch = Inches(1.5); y0 = Inches(2.55)
for i, (t, d, acc) in enumerate(news):
    col = i % 3; row = i // 3
    x = x0 + col*(cw+gap); y = y0 + row*(ch+Inches(0.22))
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, cw, Inches(0.09), acc)
    txt(s, x+Inches(0.25), y+Inches(0.25), cw-Inches(0.45), Inches(1.1),
        [[R("✦ ", 14, acc, True), R(t, 15, WHITE, True)], [R(d, 12, MUTED)]],
        space_after=5, line_spacing=1.05)
footer(s, 12, "Growth")

# =========================================================================
# SLIDE 13 — ROADMAP
# =========================================================================
s = slide(); bg(s)
header(s, "10  ·  What's Next", "The Road Ahead", AMBER)
road = [
    ("FINISH EVERY JOB", "Complete the Chef and Firefighter mini-games to full depth.", AMBER),
    ("CAREER QUEST MODE", "A progression path that guides players through all four jobs.", CYAN),
    ("MORE CAREERS", "Our framework makes adding the next job straightforward.", GREEN),
    ("DEEPER LEARNING", "In-game tips that explain the real concept behind each task.", BLUE),
    ("FRIENDS & PARTIES", "Invite friends directly and form persistent teams.", CYAN),
    ("ACCESSIBILITY & POLISH", "Tutorials, settings, and quality-of-life across the board.", GREEN),
]
x0 = Inches(0.78); cw = Inches(5.85); gap = Inches(0.35); ch = Inches(1.18); y0 = Inches(2.05)
for i, (t, d, acc) in enumerate(road):
    col = i // 3; row = i % 3
    x = x0 + col*(cw+gap); y = y0 + row*(ch+Inches(0.2))
    rect(s, x, y, cw, ch, NAVY2)
    rect(s, x, y, Inches(0.09), ch, acc)
    txt(s, x+Inches(0.3), y+Inches(0.18), cw-Inches(0.55), Inches(0.95),
        [[R(t, 15, WHITE, True)], [R(d, 12, MUTED)]], space_after=3, line_spacing=1.05)
footer(s, 13, "What's Next")

# =========================================================================
# SLIDE 14 — THANK YOU
# =========================================================================
s = slide(); bg(s)
rect(s, 0, 0, SW, Inches(0.18), CYAN)
rect(s, 0, Inches(7.32), SW, Inches(0.18), CYAN)
txt(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.3),
    [[R("Thank You", 64, WHITE, True)]])
txt(s, Inches(0.95), Inches(3.5), Inches(11), Inches(0.6),
    [[R("We'd love to answer your questions about Clock-In.", 20, CYAN, it=True)]])
rect(s, Inches(0.95), Inches(4.25), Inches(4.6), Inches(0.04), CYAN)
txt(s, Inches(0.95), Inches(4.5), Inches(11), Inches(1.2),
    [[R("CLOCK-IN", 16, BLUE, True), R("  ·  Parallel Play", 16, MUTED)],
     [R("Jackson Krizan   ·   Blayne Campbell   ·   Jakob Escueta", 17, LIGHT, True)],
     [R("FBLA National Leadership Conference 2026  ·  Computer Game & Simulation Programming", 13, MUTED)]],
    space_after=6)

prs.save("/mnt/c/Users/jacks/OneDrive/Documents/GitHub/Clock-In_V2/Docs/Clock-In - FBLA Nationals 2026.pptx")
print("Saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
